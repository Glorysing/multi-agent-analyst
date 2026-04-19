"""
一键导出 PPT
=============
从分析结果 (Markdown 报告 + 图表文件列表) 生成一份可直接转发的 PPTX。

结构:
  封面   —— 深色背景, 大标题, 目标, 日期
  章节页 —— 根据 Markdown 里的二级标题 (## 背景 / ## 关键发现 / ...) 各一张
  图表页 —— 每张图一张, 图 + 文件名小字
  结尾页 —— 深色背景 + "本 PPT 由业务数据分析系统自动生成"

原则:
  - 颜色克制: 一个主色 + 一个浅底 + 一个深底, 不乱加装饰性色条
  - 文字受控: 每章节拆成短 bullet, 单页不超过 ~8 条
  - 图片保真: fit 到 safe area, 不拉伸

依赖: python-pptx (requirements.txt 已加)
"""

from __future__ import annotations
import re
import datetime as _dt
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn


# ---- 配色 (Midnight Executive, 沉稳商务) ----
COLOR_BG_DARK   = RGBColor(0x1E, 0x27, 0x61)   # 深蓝封面/尾页
COLOR_BG_LIGHT  = RGBColor(0xFF, 0xFF, 0xFF)   # 白色正文
COLOR_ACCENT    = RGBColor(0x38, 0xBD, 0xF8)   # 天青色强调
COLOR_TEXT_DARK = RGBColor(0x0F, 0x17, 0x2A)   # 正文深色
COLOR_TEXT_MUTED = RGBColor(0x64, 0x74, 0x8B)  # 辅助
COLOR_TEXT_LIGHT = RGBColor(0xFF, 0xFF, 0xFF)  # 封面白字

# ---- 字体 ----
FONT_TITLE = "Calibri"
FONT_BODY  = "Calibri"

# ---- 页面几何 (16:9) ----
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN  = Inches(0.6)


# ----------------------------------------------------------------------
# Markdown 报告解析
# ----------------------------------------------------------------------

def _split_sections(md: str) -> list[tuple[str, list[str]]]:
    """
    把 Markdown 按 `## 标题` 切成章节.
    返回 [(标题, [bullets]), ...].
    bullet 从段落 / 有序列表 / 无序列表里提取, 每条去掉前缀符号。
    一章 bullet 太多时自动截到 8 条避免单页溢出。
    """
    if not md:
        return []
    sections: list[tuple[str, list[str]]] = []
    current_title: str | None = None
    current_lines: list[str] = []

    def flush():
        if current_title is None:
            return
        bullets = _extract_bullets(current_lines)
        if bullets:
            sections.append((current_title, bullets[:8]))

    for raw in md.splitlines():
        line = raw.rstrip()
        # 只认二级标题; 一级标题太大不适合做章节页
        m = re.match(r"^##\s+(.+?)\s*$", line)
        if m:
            flush()
            current_title = m.group(1).strip("#*_ ").strip()
            current_lines = []
        else:
            current_lines.append(line)
    flush()
    return sections


def _extract_bullets(lines: list[str]) -> list[str]:
    """从段落文本里抽 bullet. 支持 `1.` / `-` / `*` / 纯段落."""
    bullets: list[str] = []
    buf: list[str] = []

    def push_buf():
        if not buf:
            return
        text = " ".join(s.strip() for s in buf if s.strip())
        text = _strip_md(text)
        if text:
            bullets.append(text)
        buf.clear()

    for ln in lines:
        s = ln.strip()
        if not s:
            push_buf()
            continue
        m = re.match(r"^(?:\d+\.|[-*])\s+(.*)$", s)
        if m:
            push_buf()
            buf.append(m.group(1))
            push_buf()
        else:
            buf.append(s)
    push_buf()
    return bullets


def _strip_md(s: str) -> str:
    """剥掉 Markdown 强调符, 保留纯文本."""
    s = re.sub(r"\*\*(.+?)\*\*", r"\1", s)
    s = re.sub(r"__(.+?)__", r"\1", s)
    s = re.sub(r"(?<!\*)\*(?!\*)(.+?)\*(?!\*)", r"\1", s)
    s = re.sub(r"`([^`]+)`", r"\1", s)
    return s.strip()


# ----------------------------------------------------------------------
# 底层画页辅助
# ----------------------------------------------------------------------

def _set_solid_fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _add_rect(slide, x, y, w, h, color: RGBColor):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    _set_solid_fill(shape, color)
    return shape


def _add_textbox(slide, x, y, w, h, text: str, *,
                 font_size: int = 16,
                 bold: bool = False,
                 color: RGBColor = COLOR_TEXT_DARK,
                 font_name: str = FONT_BODY,
                 align=PP_ALIGN.LEFT,
                 anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)

    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def _add_bulleted_textbox(slide, x, y, w, h, bullets: list[str], *,
                          font_size: int = 16,
                          color: RGBColor = COLOR_TEXT_DARK,
                          accent: RGBColor = COLOR_ACCENT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)

    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(10)

        r_dot = p.add_run()
        r_dot.text = "• "
        r_dot.font.name = FONT_BODY
        r_dot.font.size = Pt(font_size)
        r_dot.font.bold = True
        r_dot.font.color.rgb = accent

        r_txt = p.add_run()
        r_txt.text = text
        r_txt.font.name = FONT_BODY
        r_txt.font.size = Pt(font_size)
        r_txt.font.color.rgb = color
    return tb


def _fit_image_dims(img_path: Path, max_w: int, max_h: int) -> tuple[int, int]:
    """
    读图片 (EMU 单位) 等比缩放到 (max_w, max_h) 盒子内, 返回 (w, h).
    不强转 PIL 依赖: 用 python-pptx 本身的 add_picture 时, 先传 w/h=None 拿原始尺寸。
    这里用 PIL 更稳, 但避免新增硬依赖: 退化成直接 add_picture + 自适应。
    """
    # pptx 的 add_picture 支持只给 width 或只给 height, 会自动按原比例补另一个维度。
    # 所以这里策略是: 先两次都给, 按 max 来 clip。调用方其实只用到这两个值当上限。
    return max_w, max_h


# ----------------------------------------------------------------------
# 页面构造
# ----------------------------------------------------------------------

def _blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])  # Blank


def _add_cover(prs: Presentation, title: str, subtitle: str, date_str: str):
    slide = _blank_slide(prs)
    # 深底
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, COLOR_BG_DARK)
    # 右上装饰: 一个强调色的小圆
    dot = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, SLIDE_W - Inches(1.4), Inches(0.8),
        Inches(0.6), Inches(0.6),
    )
    _set_solid_fill(dot, COLOR_ACCENT)

    # 标题区
    _add_textbox(
        slide, MARGIN, Inches(2.6), SLIDE_W - 2 * MARGIN, Inches(1.6),
        title, font_size=48, bold=True, color=COLOR_TEXT_LIGHT,
        font_name=FONT_TITLE,
    )
    # 副标题 / 分析目标
    _add_textbox(
        slide, MARGIN, Inches(4.3), SLIDE_W - 2 * MARGIN, Inches(1.3),
        subtitle, font_size=20, color=RGBColor(0xCA, 0xDC, 0xFC),
    )
    # 日期脚注
    _add_textbox(
        slide, MARGIN, SLIDE_H - Inches(0.9),
        SLIDE_W - 2 * MARGIN, Inches(0.5),
        date_str, font_size=12, color=RGBColor(0xCA, 0xDC, 0xFC),
    )


def _add_section_slide(prs: Presentation, section_title: str, bullets: list[str],
                       page_tag: str = ""):
    slide = _blank_slide(prs)
    # 白底 (默认就是), 留一个强调色左侧窄条
    _add_rect(slide, 0, 0, Inches(0.25), SLIDE_H, COLOR_ACCENT)

    # 顶部标题
    _add_textbox(
        slide, Inches(0.9), Inches(0.55),
        SLIDE_W - Inches(0.9) - MARGIN, Inches(1.0),
        section_title, font_size=32, bold=True,
        color=COLOR_TEXT_DARK, font_name=FONT_TITLE,
    )

    # 正文 bullet
    _add_bulleted_textbox(
        slide, Inches(0.9), Inches(1.8),
        SLIDE_W - Inches(0.9) - MARGIN, Inches(5.0),
        bullets, font_size=18, color=COLOR_TEXT_DARK,
    )

    # 右下角小标签 (可选)
    if page_tag:
        _add_textbox(
            slide, SLIDE_W - Inches(3.2), SLIDE_H - Inches(0.55),
            Inches(3.0), Inches(0.4),
            page_tag, font_size=10, color=COLOR_TEXT_MUTED,
            align=PP_ALIGN.RIGHT,
        )


def _add_chart_slide(prs: Presentation, chart_path: Path, caption: str = ""):
    slide = _blank_slide(prs)
    _add_rect(slide, 0, 0, Inches(0.25), SLIDE_H, COLOR_ACCENT)

    # 标题
    title = caption or chart_path.stem.replace("_", " ")
    _add_textbox(
        slide, Inches(0.9), Inches(0.45),
        SLIDE_W - Inches(0.9) - MARGIN, Inches(0.8),
        title, font_size=22, bold=True,
        color=COLOR_TEXT_DARK, font_name=FONT_TITLE,
    )

    # 图片居中贴到正文区, 最大 w * h 盒子
    box_x = Inches(0.9)
    box_y = Inches(1.4)
    box_w = SLIDE_W - Inches(0.9) - MARGIN
    box_h = Inches(5.5)

    # 直接让 pptx 按 width 缩放, 再看高度是否超, 超了就改按 height 缩放
    pic = slide.shapes.add_picture(str(chart_path), box_x, box_y, width=box_w)
    if pic.height > box_h:
        # 清掉刚加的, 改按 height 放
        sp = pic._element
        sp.getparent().remove(sp)
        pic = slide.shapes.add_picture(str(chart_path), box_x, box_y, height=box_h)
    # 居中
    pic.left = box_x + (box_w - pic.width) // 2
    pic.top = box_y + (box_h - pic.height) // 2


def _add_closing(prs: Presentation, footer_text: str):
    slide = _blank_slide(prs)
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, COLOR_BG_DARK)

    _add_textbox(
        slide, MARGIN, Inches(3.0), SLIDE_W - 2 * MARGIN, Inches(1.4),
        "Thank You", font_size=54, bold=True,
        color=COLOR_TEXT_LIGHT, font_name=FONT_TITLE, align=PP_ALIGN.CENTER,
    )
    _add_textbox(
        slide, MARGIN, Inches(4.6), SLIDE_W - 2 * MARGIN, Inches(1.0),
        footer_text, font_size=16, color=RGBColor(0xCA, 0xDC, 0xFC),
        align=PP_ALIGN.CENTER,
    )


# ----------------------------------------------------------------------
# 对外入口
# ----------------------------------------------------------------------

def build_report_pptx(
    *,
    title: str,
    user_goal: str,
    report_md: str,
    chart_paths: list[Path | str],
    language: str,
    out_path: Path | str,
) -> Path:
    """
    生成 PPTX 到 out_path. 返回 out_path (pathlib.Path).

    title        : 封面大标题 (e.g., 业务数据分析报告)
    user_goal    : 封面副标题 (展示用户分析目标)
    report_md    : Reporter 产出的 Markdown 报告
    chart_paths  : 图表绝对路径列表
    language     : "zh" | "en", 影响尾页/日期格式
    out_path     : 输出 .pptx 路径
    """
    out = Path(out_path)
    out.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    today = _dt.date.today()
    if language == "en":
        date_str = today.strftime("%B %d, %Y")
        closing = "Generated by Business Data Analyst"
        cover_title_fallback = "Business Data Analysis Report"
        charts_section_title = "Charts"
    else:
        date_str = f"{today.year}年{today.month}月{today.day}日"
        closing = "本 PPT 由业务数据分析系统自动生成"
        cover_title_fallback = "业务数据分析报告"
        charts_section_title = "图表"

    # 1. 封面
    _add_cover(
        prs,
        title=title or cover_title_fallback,
        subtitle=user_goal or "",
        date_str=date_str,
    )

    # 2. 报告各章节
    sections = _split_sections(report_md)
    for sec_title, bullets in sections:
        _add_section_slide(prs, sec_title, bullets)

    # 3. 图表
    valid_charts = [Path(p) for p in chart_paths if p and Path(p).exists()]
    if valid_charts:
        # 用一张过渡页引导 "接下来是图表"
        _add_section_slide(
            prs,
            charts_section_title,
            [f"{Path(p).stem}" for p in valid_charts],
            page_tag=f"{len(valid_charts)} / charts",
        )
        for p in valid_charts:
            _add_chart_slide(prs, p)

    # 4. 结尾
    _add_closing(prs, closing)

    prs.save(str(out))
    return out
